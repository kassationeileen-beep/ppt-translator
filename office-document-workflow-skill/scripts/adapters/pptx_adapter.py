"""PPTX extraction/writeback adapter using python-pptx."""
from __future__ import annotations

from pathlib import Path
from typing import Dict, Iterable, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from scripts.core.document_model import LayoutMetadata, RiskFlags, StyleMetadata, TextBlock
from scripts.core.qa import likely_pptx_overflow
from scripts.core.text_spacing import add_run_boundary_space_if_needed


def _color_to_hex(run) -> str | None:
    color = run.font.color
    if color is None or color.rgb is None:
        return None
    return str(color.rgb)


def _font_size_points(run) -> float | None:
    return float(run.font.size.pt) if run.font.size is not None else None


PPTX_MIN_FONT_SIZE_PT = 6
PPTX_DEFAULT_FONT_SIZE_PT = 12
PPTX_OVERFLOW_SCALE_STEP = 0.85
PPTX_MAX_FIT_PASSES = 4


class PptxAdapter:
    document_type = "pptx"

    def __init__(self, path: str | Path):
        self.path = Path(path)
        self.presentation = Presentation(str(self.path))

    def extract_text_blocks(self) -> List[TextBlock]:
        blocks: List[TextBlock] = []
        for slide_index, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                blocks.extend(self._extract_shape(shape, slide_index))
        return blocks

    def _extract_shape(self, shape, slide_index: int) -> List[TextBlock]:
        blocks: List[TextBlock] = []
        shape_id = getattr(shape, "shape_id", None)
        layout = LayoutMetadata(
            x=getattr(shape, "left", None),
            y=getattr(shape, "top", None),
            width=getattr(shape, "width", None),
            height=getattr(shape, "height", None),
        )
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                child_blocks = self._extract_shape(child, slide_index)
                for block in child_blocks:
                    block.risks.is_grouped_shape = True
                    block.risks.requires_manual_review = True
                blocks.extend(child_blocks)
            return blocks
        if getattr(shape, "has_chart", False):
            blocks.append(TextBlock(
                document_type="pptx",
                block_id=f"pptx:s{slide_index}:shape{shape_id}:chart",
                slide_index=slide_index,
                shape_id=shape_id,
                paragraph_index=None,
                original_text="",
                layout=layout,
                risks=RiskFlags(is_chart_text=True, requires_manual_review=True),
                context={"kind": "chart"},
            ))
            return blocks
        if getattr(shape, "has_table", False):
            for row_index, row in enumerate(shape.table.rows):
                for col_index, cell in enumerate(row.cells):
                    for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs):
                        blocks.extend(self._paragraph_blocks(
                            paragraph,
                            slide_index,
                            shape_id,
                            paragraph_index,
                            layout,
                            risks=RiskFlags(is_table_text=True),
                            context={"kind": "table", "row_index": row_index, "col_index": col_index},
                        ))
            return blocks
        if getattr(shape, "has_text_frame", False):
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                blocks.extend(self._paragraph_blocks(
                    paragraph,
                    slide_index,
                    shape_id,
                    paragraph_index,
                    layout,
                    risks=RiskFlags(),
                    context={"kind": "text_frame"},
                ))
            return blocks
        if getattr(shape, "shape_type", None) not in {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER, MSO_SHAPE_TYPE.AUTO_SHAPE}:
            blocks.append(TextBlock(
                document_type="pptx",
                block_id=f"pptx:s{slide_index}:shape{shape_id}:unsupported",
                slide_index=slide_index,
                shape_id=shape_id,
                paragraph_index=None,
                original_text="",
                layout=layout,
                risks=RiskFlags(is_smartart_or_unsupported=True, requires_manual_review=True),
                context={"kind": "unsupported_shape", "shape_type": str(getattr(shape, "shape_type", None))},
            ))
        return blocks

    def _paragraph_blocks(self, paragraph, slide_index, shape_id, paragraph_index, layout, risks, context) -> List[TextBlock]:
        blocks: List[TextBlock] = []
        alignment = str(paragraph.alignment) if paragraph.alignment is not None else None
        bullet_level = getattr(paragraph, "level", None)
        for run_index, run in enumerate(paragraph.runs):
            if not run.text:
                continue
            block_id = f"pptx:s{slide_index}:shape{shape_id}:p{paragraph_index}:r{run_index}"
            blocks.append(TextBlock(
                document_type="pptx",
                block_id=block_id,
                slide_index=slide_index,
                shape_id=shape_id,
                paragraph_index=paragraph_index,
                run_index=run_index,
                original_text=run.text,
                style=StyleMetadata(
                    font_name=run.font.name,
                    font_size=_font_size_points(run),
                    bold=run.font.bold,
                    italic=run.font.italic,
                    underline=run.font.underline,
                    color=_color_to_hex(run),
                    alignment=alignment,
                    bullet_level=bullet_level,
                ),
                layout=layout,
                risks=risks,
                context=context,
            ))
        if not blocks and paragraph.text:
            block_id = f"pptx:s{slide_index}:shape{shape_id}:p{paragraph_index}:text"
            blocks.append(TextBlock(
                document_type="pptx",
                block_id=block_id,
                slide_index=slide_index,
                shape_id=shape_id,
                paragraph_index=paragraph_index,
                original_text=paragraph.text,
                style=StyleMetadata(alignment=alignment, bullet_level=bullet_level),
                layout=layout,
                risks=risks,
                context={**context, "paragraph_level": True},
            ))
        return blocks

    def write_translations(self, blocks: Iterable[TextBlock], output_path: str | Path) -> None:
        block_map: Dict[str, TextBlock] = {block.block_id: block for block in blocks if block.translated_text and not block.risks.requires_manual_review}
        for slide_index, slide in enumerate(self.presentation.slides):
            for shape in slide.shapes:
                self._write_shape(shape, slide_index, block_map)
        self.presentation.save(str(output_path))

    def _write_shape(self, shape, slide_index: int, block_map: Dict[str, TextBlock]) -> None:
        shape_id = getattr(shape, "shape_id", None)
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                self._write_shape(child, slide_index, block_map)
            return
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                for cell in row.cells:
                    cell_may_overflow = False
                    for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs):
                        paragraph_may_overflow = self._write_paragraph(
                            paragraph,
                            slide_index,
                            shape_id,
                            paragraph_index,
                            block_map,
                        )
                        cell_may_overflow = paragraph_may_overflow or cell_may_overflow
                    if cell_may_overflow:
                        self._enable_text_frame_autofit(cell.text_frame)
            return
        if getattr(shape, "has_text_frame", False):
            shape_may_overflow = False
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                paragraph_may_overflow = self._write_paragraph(
                    paragraph,
                    slide_index,
                    shape_id,
                    paragraph_index,
                    block_map,
                )
                shape_may_overflow = paragraph_may_overflow or shape_may_overflow
            if shape_may_overflow:
                self._enable_text_frame_autofit(shape.text_frame)

    def _write_paragraph(
        self,
        paragraph,
        slide_index: int,
        shape_id: int,
        paragraph_index: int,
        block_map: Dict[str, TextBlock],
    ) -> bool:
        paragraph_blocks: List[TextBlock] = []
        if paragraph.runs:
            previous_block: TextBlock | None = None
            previous_text = ""
            for run_index, run in enumerate(paragraph.runs):
                block = block_map.get(f"pptx:s{slide_index}:shape{shape_id}:p{paragraph_index}:r{run_index}")
                if block:
                    translated_text = block.translated_text or run.text
                    if previous_block is not None:
                        translated_text = add_run_boundary_space_if_needed(
                            previous_block.original_text,
                            block.original_text,
                            previous_text,
                            translated_text,
                        )
                        block.translated_text = translated_text
                    run.text = translated_text
                    paragraph_blocks.append(block)
                    previous_block = block
                    previous_text = translated_text
                elif run.text:
                    previous_block = None
                    previous_text = run.text
        else:
            block = block_map.get(f"pptx:s{slide_index}:shape{shape_id}:p{paragraph_index}:text")
            if block:
                paragraph.text = block.translated_text or paragraph.text
                paragraph_blocks.append(block)

        if not paragraph_blocks:
            return False
        may_overflow = any(likely_pptx_overflow(block) for block in paragraph_blocks) or self._paragraph_expanded(
            paragraph_blocks
        )
        if may_overflow:
            for block in paragraph_blocks:
                block.risks.may_overflow = True
            self._shrink_paragraph_runs(paragraph)
        return may_overflow

    @staticmethod
    def _paragraph_expanded(blocks: List[TextBlock]) -> bool:
        original_len = sum(len(block.original_text or "") for block in blocks)
        translated_len = sum(len(block.translated_text or "") for block in blocks)
        return original_len > 0 and translated_len / original_len > 1.45

    @staticmethod
    def _enable_text_frame_autofit(text_frame) -> None:
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    @staticmethod
    def _shrink_paragraph_runs(paragraph) -> None:
        for _ in range(PPTX_MAX_FIT_PASSES):
            changed = False
            for run in paragraph.runs:
                current_size = run.font.size.pt if run.font.size is not None else PPTX_DEFAULT_FONT_SIZE_PT
                next_size = max(current_size * PPTX_OVERFLOW_SCALE_STEP, PPTX_MIN_FONT_SIZE_PT)
                if next_size < current_size:
                    run.font.size = Pt(next_size)
                    changed = True
            if not changed:
                break
